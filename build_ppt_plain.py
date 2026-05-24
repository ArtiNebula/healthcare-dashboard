"""Build 21-slide HealthAI presentation — NO theme, plain black text on white."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT  = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\HealthAI_Plain.pptx'
SCRN = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\screenshots'
IMGS = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\dissertation_imgs'

BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY  = RGBColor(0x44, 0x44, 0x44)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # Blank

# ── helpers ───────────────────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(blank)

def title(sl, text, t=0.25, size=32):
    box = sl.shapes.add_textbox(Inches(0.5), Inches(t), Inches(12.3), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLACK

def subtitle(sl, text, t=0.95, size=14):
    box = sl.shapes.add_textbox(Inches(0.5), Inches(t), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = GRAY

def body(sl, lines, l=0.5, t=1.45, w=12.3, h=5.5, size=13):
    """lines = list of (indent_level, text_str)  indent 0=bullet 1=sub"""
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (lvl, txt) in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(2)
        bullet_char = '•' if lvl == 0 else '–'
        run = p.add_run()
        run.text = f'{bullet_char}  {txt}'
        run.font.size = Pt(size - lvl * 1.5)
        run.font.color.rgb = BLACK

def txt(sl, text, l, t, w, h, size=12, bold=False, align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = BLACK

def img(sl, path, l, t, w, h=None):
    if not os.path.exists(path):
        print(f'  WARN: {path}')
        return
    if h:
        sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        sl.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))

def caption(sl, text, l, t, w):
    txt(sl, text, l, t, w, 0.3, size=9, italic=True, align=PP_ALIGN.CENTER)

def col_section(sl, col_x, col_t, col_title, items, col_w=3.9, item_size=11):
    txt(sl, col_title, col_x, col_t, col_w, 0.4, size=14, bold=True)
    for j, item in enumerate(items):
        txt(sl, f'•  {item}', col_x, col_t + 0.45 + j*0.38, col_w, 0.35, size=item_size)

def two_col_body(sl, left_lines, right_lines, t=1.45, h=5.5, size=12):
    body(sl, left_lines,  l=0.5, t=t, w=6.0, h=h, size=size)
    body(sl, right_lines, l=6.8, t=t, w=6.0, h=h, size=size)

def pagenum(sl, num, total=21):
    txt(sl, f'{num} / {total}', 12.5, 7.2, 0.7, 0.25, size=9, align=PP_ALIGN.RIGHT)

def hline(sl, t):
    from pptx.util import Emu
    ln = sl.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.3), Emu(18000))
    ln.fill.background()
    ln.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    ln.line.width = Pt(0.5)

# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
sl = slide()
txt(sl, 'HealthAI', 0, 1.5, 13.333, 1.2, size=54, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'AI-Powered Health Symptom Tracker', 0, 2.75, 13.333, 0.55,
    size=20, align=PP_ALIGN.CENTER)
txt(sl, 'Deployed on Cloud-Native Infrastructure using Docker & Kubernetes (K3s)',
    1, 3.35, 11.333, 0.45, size=13, align=PP_ALIGN.CENTER)
hline(sl, 4.1)
txt(sl, 'Arti Kumari  |  2024MT03117  |  M.Tech Cloud Computing',
    0, 4.25, 13.333, 0.4, size=14, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'BITS Pilani — Work Integrated Learning Programmes  |  May 2026',
    0, 4.68, 13.333, 0.35, size=12, align=PP_ALIGN.CENTER)
txt(sl, 'Supervisor: Atul Kumar (Solution Architect), ValueLabs, Indore',
    0, 5.08, 13.333, 0.35, size=11, align=PP_ALIGN.CENTER)
pagenum(sl, 1)
print('Slide 1: Title')

# =============================================================================
# SLIDE 2 — THE PROBLEM
# =============================================================================
sl = slide()
title(sl, 'Why This Matters — The Problem Space in Digital Healthcare')
subtitle(sl, 'Three systemic gaps that motivated this project')
hline(sl, 1.35)
body(sl, [
    (0, 'Reactive, Not Proactive'),
    (1, 'Patients log symptoms but receive no AI-driven insights or early warnings'),
    (1, 'No confidence scoring, trend analysis, or severity-based escalation'),
    (0, 'Monolithic, Not Scalable'),
    (1, 'Single-server apps collapse under traffic spikes — no elasticity or auto-recovery'),
    (1, 'Deployment involves manual SSH, no rollback, no canary releases'),
    (0, 'Observable on Paper, Blind in Practice'),
    (1, 'No real-time dashboards — IT teams discover outages from angry users'),
    (1, 'Zero custom metrics, no alerting, no distributed tracing'),
    (0, 'The Stakes'),
    (1, '1 in 3 patients: "I didn\'t know when to see a doctor" (WHO 2022)'),
    (1, 'Average healthcare app downtime: 14 hours / year'),
    (1, 'HIPAA violations: avg fine $1.19M per incident (2023)'),
], size=12)
pagenum(sl, 2)
print('Slide 2: Problem')

# =============================================================================
# SLIDE 3 — SOLUTION
# =============================================================================
sl = slide()
title(sl, 'Introducing HealthAI — One Platform, Three Capabilities')
subtitle(sl, 'What we built and why it matters')
hline(sl, 1.35)

col_section(sl, 0.4, 1.45, 'Patient Interface', [
    'Log symptoms in under 1 minute',
    'AI confidence scoring per entry',
    'Health trend charts (Line/Area/Radar)',
    '6-category symptom grid',
    'Personal health history timeline',
])
col_section(sl, 4.7, 1.45, 'Admin Interface', [
    'Monitor all patients in real-time',
    'Severe case flagging and alerts',
    'Population-level analytics',
    'User growth and engagement charts',
    'One-click alert management',
])
col_section(sl, 9.0, 1.45, 'Cloud Infrastructure', [
    'Auto-scale 2 → 6 pods via HPA',
    'Prometheus + Grafana monitoring',
    'HIPAA §164.312 compliant design',
    '0.08% error rate @ 100 virtual users',
    'Immutable Docker containers',
])

hline(sl, 5.8)
txt(sl, '8 pages  •  15+ REST endpoints  •  6 custom Prometheus metrics  •  All HIPAA §164.312 safeguards covered',
    0.5, 5.88, 12.3, 0.3, size=10, align=PP_ALIGN.CENTER)
pagenum(sl, 3)
print('Slide 3: Solution')

# =============================================================================
# SLIDE 4 — TECHNOLOGY STACK
# =============================================================================
sl = slide()
title(sl, 'Built With Production-Grade Tools')
subtitle(sl, 'Full-stack technology selection rationale')
hline(sl, 1.35)

col_section(sl, 0.4, 1.45, 'Frontend', [
    'React 18 + TypeScript — type-safe UI',
    'Vite — sub-second HMR builds',
    'Tailwind CSS v4 — utility-first styling',
    'Recharts — Line/Area/Bar/Radar charts',
    'React Router v7 — nested layouts',
    'shadcn/ui — accessible components',
], col_w=4.0)
col_section(sl, 4.7, 1.45, 'Backend', [
    'Node.js + Express.js — REST API',
    'MySQL 8.x — relational health data',
    'bcryptjs (factor 10) — password hashing',
    'ioredis — Redis cache (30s TTL)',
    'express-rate-limit — brute-force guard',
    'prom-client — Prometheus metrics',
], col_w=4.0)
col_section(sl, 9.0, 1.45, 'Infrastructure', [
    'Docker multi-stage builds',
    'Docker Compose — local orchestration',
    'K3s — lightweight Kubernetes (512MB)',
    'HPA — CPU autoscaling (2–6 pods)',
    'Traefik — ingress + TLS termination',
    'Prometheus + Grafana — observability',
], col_w=4.0)
pagenum(sl, 4)
print('Slide 4: Tech Stack')

# =============================================================================
# SLIDE 5 — SYSTEM ARCHITECTURE
# =============================================================================
sl = slide()
title(sl, 'Three-Tier Cloud-Native Architecture')
subtitle(sl, 'System design overview')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig3_1_arch.png'), 6.8, 1.45, 6.3, 5.5)

body(sl, [
    (0, 'Presentation Tier'),
    (1, 'React 18 + TypeScript SPA (Nginx container)'),
    (1, 'X-User-Id / X-User-Email auth headers'),
    (0, 'Application Tier'),
    (1, 'Node.js + Express.js REST API'),
    (1, 'RBAC enforcement + AI analysis engine'),
    (1, 'prom-client custom metrics'),
    (0, 'Data Tier'),
    (1, 'MySQL 8.x relational database'),
    (1, 'Redis cache (30-second TTL)'),
    (0, 'Infrastructure Layer'),
    (1, 'Docker containers per service'),
    (1, 'K3s orchestration + Traefik ingress'),
    (1, 'Prometheus scrape every 15s → Grafana'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 5)
print('Slide 5: Architecture')

# =============================================================================
# SLIDE 6 — FRONTEND ARCHITECTURE
# =============================================================================
sl = slide()
title(sl, 'React 18 SPA — Role-Based Route Layout')
subtitle(sl, 'Frontend structure and navigation hierarchy')
hline(sl, 1.35)

route_lines = """/                        Root Layout (public Navbar)
  /                      Home Page
  /login                 Login Page
  /signup                Signup Page

/user                    User Layout (sidebar, 5 nav items)
  /user                  Dashboard Page
  /user/add-symptoms     Add Symptoms Page
  /user/history          Health History Page
  /user/ai-suggestions   AI Suggestion Page
  /user/profile          Profile Page

/admin                   Admin Layout (admin sidebar)
  /admin                 Admin Dashboard Page
  /admin/monitoring      User Monitoring Page
  /admin/severe-cases    Severe Cases Page
  /admin/alerts          Alerts Page
  /admin/reports         Reports Page"""

box = sl.shapes.add_textbox(Inches(0.4), Inches(1.45), Inches(5.8), Inches(5.5))
box.word_wrap = True
tf = box.text_frame
tf.word_wrap = True
first = True
for line in route_lines.split('\n'):
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(9.5)
    run.font.name = 'Courier New'
    run.font.color.rgb = BLACK

body(sl, [
    (0, 'Role-Based Access Control'),
    (1, 'Patient users see /user/* routes only'),
    (1, 'Admin users see /admin/* routes only'),
    (1, 'RBAC enforced at backend and frontend'),
    (0, 'Authentication'),
    (1, 'Stateless: user object in localStorage'),
    (1, 'X-User-Id header on every API call'),
    (1, 'bcrypt password validation on login'),
    (0, 'Protected Routes'),
    (1, 'Redirect to /login if no session'),
    (1, 'Redirect to /admin-login for admin'),
], l=6.8, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 6)
print('Slide 6: Frontend')

# =============================================================================
# SLIDE 7 — USER DASHBOARD
# =============================================================================
sl = slide()
title(sl, 'Patient Dashboard — Feature Walkthrough')
subtitle(sl, 'Six core pages of the patient-facing interface')
hline(sl, 1.35)

img(sl, os.path.join(SCRN, '02_user_dashboard.png'), 0.4, 1.45, 5.9, 3.7)
caption(sl, 'Figure: Patient Dashboard — health summary & recent symptoms', 0.4, 5.2, 5.9)

body(sl, [
    (0, 'Dashboard'),
    (1, 'Stat cards: BMI, Heart Rate, Steps, Sleep'),
    (1, 'Recent symptoms with severity badges'),
    (1, 'Area chart: 7-day health trend'),
    (0, 'Add Symptoms'),
    (1, '6-category grid + free-text input'),
    (1, 'Immediate AI confidence score'),
    (0, 'Health History'),
    (1, 'Chronological symptom timeline'),
    (1, 'Filter by date/severity/category'),
    (0, 'AI Suggestions'),
    (1, 'Weighted scoring across 5 categories'),
    (1, 'Color-coded recommendations'),
], l=6.8, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 7)
print('Slide 7: Dashboard')

# =============================================================================
# SLIDE 8 — AI ANALYSIS ENGINE
# =============================================================================
sl = slide()
title(sl, 'AI Analysis Engine — Symptom Intelligence')
subtitle(sl, 'Rule-based + weighted scoring system for clinical suggestions')
hline(sl, 1.35)

img(sl, os.path.join(SCRN, '05_ai_suggestions.png'), 6.8, 1.45, 6.1, 3.7)
caption(sl, 'Figure: AI Suggestions Page', 6.8, 5.2, 6.1)

body(sl, [
    (0, 'Scoring Algorithm'),
    (1, 'Severity weights: Mild=1, Moderate=2, Severe=3'),
    (1, 'Time decay: Today=1.0, Yesterday=0.7, 2d=0.5'),
    (1, 'Category weights × recency → total risk score'),
    (0, 'Output Categories'),
    (1, 'Respiratory, Cardiovascular, Neurological'),
    (1, 'Gastrointestinal, Musculoskeletal'),
    (1, 'Severity threshold → "Seek immediate care"'),
    (0, 'API Endpoint'),
    (1, 'GET /api/ai-suggestions?userId=N'),
    (1, 'Returns top 5 recommendations + confidence %'),
    (1, 'Response < 50ms (Redis-cached per user)'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 8)
print('Slide 8: AI Analysis')

# =============================================================================
# SLIDE 9 — ADMIN INTERFACE
# =============================================================================
sl = slide()
title(sl, 'Administrative Interface — Operational Control')
subtitle(sl, 'Real-time patient monitoring, alerts, and analytics')
hline(sl, 1.35)

img(sl, os.path.join(SCRN, '06_admin_dashboard.png'), 0.4, 1.45, 5.9, 3.5)
caption(sl, 'Figure: Admin Dashboard Overview', 0.4, 5.0, 5.9)

body(sl, [
    (0, 'Admin Dashboard'),
    (1, 'Total users, active today, severe cases count'),
    (1, 'Monthly user growth bar chart'),
    (1, 'Symptom category distribution radar'),
    (0, 'User Monitoring'),
    (1, 'List of all patients with symptom counts'),
    (1, 'Last activity timestamp per user'),
    (0, 'Severe Cases'),
    (1, 'Auto-flagged: Severe severity or score > 8'),
    (1, 'Contact details, history access'),
    (0, 'Alerts & Reports'),
    (1, 'Unresolved alert queue + one-click resolve'),
    (1, 'CSV export + printable PDF reports'),
], l=6.8, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 9)
print('Slide 9: Admin')

# =============================================================================
# SLIDE 10 — DOCKER CONTAINERISATION
# =============================================================================
sl = slide()
title(sl, 'Docker Containerisation Strategy')
subtitle(sl, 'Immutable builds, multi-stage Dockerfiles, Docker Compose orchestration')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig4_1_docker.png'), 6.8, 1.45, 6.1, 4.5)
caption(sl, 'Figure: Docker architecture — multi-stage builds & service layout', 6.8, 6.0, 6.1)

body(sl, [
    (0, 'Multi-Stage Dockerfile (Frontend)'),
    (1, 'Stage 1 — builder: npm ci + vite build'),
    (1, 'Stage 2 — nginx:alpine: copy dist/ only'),
    (1, 'Final image: ~22MB (vs 800MB+ full Node)'),
    (0, 'Multi-Stage Dockerfile (Backend)'),
    (1, 'Stage 1 — deps: npm ci --omit=dev'),
    (1, 'Stage 2 — node:20-alpine: minimal runtime'),
    (1, 'Non-root user (uid 1001) for security'),
    (0, 'Docker Compose Services'),
    (1, 'frontend, backend, mysql, redis, prometheus, grafana'),
    (1, 'Shared bridge network: healthcare-net'),
    (1, 'Named volumes: mysql_data, redis_data'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 10)
print('Slide 10: Docker')

# =============================================================================
# SLIDE 11 — K3s KUBERNETES
# =============================================================================
sl = slide()
title(sl, 'K3s — Lightweight Kubernetes Deployment')
subtitle(sl, 'Production orchestration on 512MB RAM with full K8s API compatibility')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig4_2_k3s.png'), 6.8, 1.45, 6.1, 4.5)
caption(sl, 'Figure: K3s cluster — namespace, pods, services, ingress', 6.8, 6.0, 6.1)

body(sl, [
    (0, 'Why K3s?'),
    (1, 'Single 60MB binary, 512MB RAM baseline'),
    (1, 'Full K8s API — same YAML as full K8s'),
    (1, 'Built-in containerd, CoreDNS, Traefik'),
    (0, 'Namespace: healthcare'),
    (1, 'Deployments: frontend, backend, mysql, redis'),
    (1, 'Services: ClusterIP for internal, NodePort for ingress'),
    (1, 'Secrets: mysql-secret, redis-secret'),
    (0, 'Traefik Ingress'),
    (1, 'IngressRoute → frontend at /'),
    (1, 'IngressRoute → backend at /api'),
    (1, 'TLS termination via cert-manager'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 11)
print('Slide 11: K3s')

# =============================================================================
# SLIDE 12 — HPA AUTOSCALING
# =============================================================================
sl = slide()
title(sl, 'Horizontal Pod Autoscaler — Dynamic Scaling')
subtitle(sl, 'CPU-based autoscaling: 2 → 6 backend pods in under 60 seconds')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig5_3_hpa.png'), 6.8, 1.45, 6.1, 4.5)
caption(sl, 'Figure: HPA replica count over 60-minute window', 6.8, 6.0, 6.1)

body(sl, [
    (0, 'HPA Configuration'),
    (1, 'minReplicas: 2  |  maxReplicas: 6'),
    (1, 'targetCPUUtilizationPercentage: 50'),
    (1, 'scaleDown stabilization: 120 seconds'),
    (0, 'Scaling Algorithm'),
    (1, 'Metrics Server polls every 15s'),
    (1, 'desiredReplicas = ceil(current × actual/target)'),
    (1, 'Scale-up: immediate; scale-down: 120s delay'),
    (0, 'Observed Behaviour (Load Test)'),
    (1, '2 pods idle → 4 pods @ 100 VUs in 45s'),
    (1, 'CPU per pod: 45–55% under load'),
    (1, 'Scale-down to 2 within 3 min post-load'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 12)
print('Slide 12: HPA')

# =============================================================================
# SLIDE 13 — MONITORING
# =============================================================================
sl = slide()
title(sl, 'Observability Stack — Prometheus + Grafana')
subtitle(sl, 'Six custom metrics, four dashboard panels, full-stack visibility')
hline(sl, 1.35)

img(sl, os.path.join(SCRN, '11_grafana.png'), 6.8, 1.45, 6.1, 4.5)
caption(sl, 'Figure: Grafana HealthAI System Overview Dashboard', 6.8, 6.0, 6.1)

body(sl, [
    (0, 'Custom Prometheus Metrics'),
    (1, 'http_requests_total — by method/route/status'),
    (1, 'http_request_duration_seconds — histogram'),
    (1, 'symptom_submissions_total — by severity'),
    (1, 'ai_analysis_duration_seconds — latency'),
    (1, 'active_users_total — gauge'),
    (1, 'severe_cases_total — clinical alert counter'),
    (0, 'Grafana Dashboard Panels'),
    (1, 'API Request Rate (req/s) — line chart'),
    (1, 'Response Time P95 (ms) — with SLA threshold'),
    (1, 'CPU / Memory per pod — multi-series'),
    (1, 'HPA Replica Count — step chart'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 13)
print('Slide 13: Monitoring')

# =============================================================================
# SLIDE 14 — SECURITY ARCHITECTURE
# =============================================================================
sl = slide()
title(sl, 'Security Architecture — Multi-Layer Defence')
subtitle(sl, 'Security-as-Code: rate limiting, RBAC, bcrypt, network policies')
hline(sl, 1.35)

two_col_body(sl, [
    (0, 'Authentication'),
    (1, 'bcrypt factor 10 — ~120ms per hash'),
    (1, 'Brute-force: 10 req/15min/IP limit'),
    (1, 'Stateless sessions via headers'),
    (0, 'Authorisation'),
    (1, 'RBAC: patient vs admin roles'),
    (1, 'Backend validates role on every request'),
    (1, 'Admin routes return 403 for patients'),
    (0, 'Input Validation'),
    (1, 'Parameterised SQL — no injection'),
    (1, 'Joi schema validation on all bodies'),
], [
    (0, 'Network Security'),
    (1, 'K8s NetworkPolicy: pod-to-pod whitelist'),
    (1, 'TLS 1.2/1.3 via Traefik + cert-manager'),
    (1, 'HTTPS enforced in production'),
    (0, 'Secrets Management'),
    (1, 'K8s Secrets: base64-encoded credentials'),
    (1, 'Docker env vars — never baked into image'),
    (0, 'Runtime Security'),
    (1, 'Non-root containers (uid 1001)'),
    (1, 'Read-only root filesystem'),
    (1, 'Trivy scan: 0 critical CVEs'),
], size=11.5)
pagenum(sl, 14)
print('Slide 14: Security')

# =============================================================================
# SLIDE 15 — HIPAA COMPLIANCE
# =============================================================================
sl = slide()
title(sl, 'HIPAA §164.312 — Technical Safeguards Coverage')
subtitle(sl, 'All seven required technical safeguard standards addressed')
hline(sl, 1.35)

# Table as text (simple)
headers = ['HIPAA §164.312 Standard', 'Requirement', 'Implementation']
rows = [
    ('(a)(1) Access Control',        'Unique user IDs, emergency access', 'RBAC roles, bcrypt auth'),
    ('(a)(2)(i) Unique User ID',     'Identify each user uniquely',       'user.id + email headers'),
    ('(a)(2)(iii) Auto Log-off',     'Terminate idle sessions',           'localStorage + expiry'),
    ('(a)(2)(iv) Encryption/Decrypt','Encrypt ePHI at rest',              'MySQL encrypted volume'),
    ('(b) Audit Controls',           'Hardware/software activity logs',   'Prometheus + prom-client'),
    ('(c) Integrity',                'Protect ePHI from alteration',      'Parameterised SQL, HTTPS'),
    ('(d) Auth / Person',            'Verify person seeking access',      'bcrypt + RBAC validation'),
    ('(e) Transmission Security',    'Guard ePHI in transit',             'TLS 1.2/1.3, HTTPS only'),
]

# Draw table as text boxes
col_ws = [3.5, 3.5, 4.8]
col_xs = [0.4, 3.95, 7.5]
y_start = 1.5
row_h  = 0.52

# Header
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
    txt(sl, hdr, cx, y_start, cw, row_h-0.05, size=10.5, bold=True)

for ri, row in enumerate(rows):
    y = y_start + row_h*(ri+1)
    for ci, (cell, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        txt(sl, cell, cx, y, cw, row_h-0.05, size=9.5)

pagenum(sl, 15)
print('Slide 15: HIPAA')

# =============================================================================
# SLIDE 16 — LOAD TESTING
# =============================================================================
sl = slide()
title(sl, 'Load Testing — k6 Performance Benchmarks')
subtitle(sl, '100 virtual users, 3-minute sustained load, 4 critical scenarios')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig5_1_response.png'), 6.8, 1.45, 6.1, 4.5)
caption(sl, 'Figure: API Response Time Distribution (P50/P95/P99)', 6.8, 6.0, 6.1)

body(sl, [
    (0, 'Test Configuration'),
    (1, 'Tool: k6 v0.48 — JS-based load scenarios'),
    (1, '100 VUs ramp-up over 30s, hold 3 min'),
    (1, 'SLO: P95 < 500ms, error rate < 1%'),
    (0, 'Scenarios Tested'),
    (1, 'S1: Patient login → dashboard → symptom add'),
    (1, 'S2: AI analysis fetch (Redis cold vs warm)'),
    (1, 'S3: Admin monitoring page under concurrent load'),
    (1, 'S4: HPA scale-out: 2 → 4 → 6 pods'),
    (0, 'Results'),
    (1, 'P50 = 21ms  |  P95 = 38ms  |  P99 = 67ms'),
    (1, 'Error rate: 0.08% (all 503 during HPA scale-up)'),
    (1, 'Throughput: 142 req/s sustained'),
], l=0.4, t=1.45, w=6.1, h=5.5, size=11.5)
pagenum(sl, 16)
print('Slide 16: Load Testing')

# =============================================================================
# SLIDE 17 — PERFORMANCE RESULTS
# =============================================================================
sl = slide()
title(sl, 'Performance Results — System Under Load')
subtitle(sl, 'CPU, memory, and response-time metrics during k6 benchmark')
hline(sl, 1.35)

img(sl, os.path.join(IMGS, 'fig5_2_cpu_mem.png'), 0.4, 1.45, 6.1, 4.0)
caption(sl, 'Figure: CPU & Memory utilisation per pod (100 VU load)', 0.4, 5.5, 6.1)

img(sl, os.path.join(IMGS, 'fig5_3_hpa.png'), 6.8, 1.45, 6.1, 4.0)
caption(sl, 'Figure: HPA replica scaling timeline', 6.8, 5.5, 6.1)

body(sl, [
    (0, 'Peak CPU: 52% per pod (3-pod average)'),
    (0, 'Peak Memory: 314 MiB per pod'),
    (0, 'HPA scaled 2 → 4 pods within 45 seconds'),
    (0, 'Scale-down: 2 pods within 3 min post-load'),
    (0, 'Redis cache-hit rate: 94% during test'),
], l=0.4, t=5.75, w=12.3, h=1.5, size=11)
pagenum(sl, 17)
print('Slide 17: Performance')

# =============================================================================
# SLIDE 18 — SECURITY SCAN RESULTS
# =============================================================================
sl = slide()
title(sl, 'Security Scan Results — Trivy + Manual Testing')
subtitle(sl, 'Zero critical CVEs in final images; OWASP Top 10 verified')
hline(sl, 1.35)

two_col_body(sl, [
    (0, 'Trivy Container Scan'),
    (1, 'frontend image (nginx:alpine) — 0 CRITICAL'),
    (1, 'backend image (node:20-alpine) — 0 CRITICAL'),
    (1, 'mysql:8 image — 2 MEDIUM (upstream, patched)'),
    (0, 'OWASP Top 10 Checks'),
    (1, 'A01 Broken Access Control — RBAC enforced'),
    (1, 'A02 Cryptographic Failures — TLS + bcrypt'),
    (1, 'A03 Injection — parameterised SQL throughout'),
    (1, 'A07 Auth Failures — rate-limit + bcrypt'),
], [
    (0, 'Penetration Test Findings'),
    (1, 'SQL injection: Not possible (parameterised)'),
    (1, 'XSS: React escapes all output by default'),
    (1, 'CSRF: Stateless API + CORS whitelist'),
    (1, 'Brute force: 429 after 10 req/15min'),
    (0, 'Kubernetes Security'),
    (1, 'NetworkPolicy: pod-to-pod whitelist'),
    (1, 'Non-root containers across all pods'),
    (1, 'Secrets encrypted at rest (etcd)'),
], size=11.5)
pagenum(sl, 18)
print('Slide 18: Security Scan')

# =============================================================================
# SLIDE 19 — CONCLUSIONS
# =============================================================================
sl = slide()
title(sl, 'Conclusions — What Was Achieved')
subtitle(sl, 'Objectives met across all four evaluation dimensions')
hline(sl, 1.35)

two_col_body(sl, [
    (0, 'Technical Achievements'),
    (1, '15+ REST endpoints, all functional'),
    (1, 'AI analysis with confidence scoring'),
    (1, 'HPA scales 2 → 6 pods in < 60s'),
    (1, 'P95 latency 38ms @ 100 VUs'),
    (1, '0.08% error rate under full load'),
    (0, 'Security & Compliance'),
    (1, 'All HIPAA §164.312 safeguards covered'),
    (1, '0 critical CVEs in container images'),
    (1, 'OWASP Top 10 verified'),
    (1, 'TLS + RBAC + bcrypt + parameterised SQL'),
], [
    (0, 'Observability'),
    (1, '6 custom Prometheus metrics'),
    (1, 'Grafana dashboard with 5+ panels'),
    (1, 'HPA step chart, CPU/memory trends'),
    (1, 'Alert manager integration ready'),
    (0, 'Project Impact'),
    (1, 'Demonstrates cloud-native healthcare at scale'),
    (1, 'Reproducible: Docker Compose one-command setup'),
    (1, 'K3s deployment on 512MB RAM VM'),
    (1, 'Complete audit trail via Prometheus'),
], size=11.5)
pagenum(sl, 19)
print('Slide 19: Conclusions')

# =============================================================================
# SLIDE 20 — FUTURE WORK
# =============================================================================
sl = slide()
title(sl, 'Future Work — Roadmap')
subtitle(sl, 'Six enhancements planned for next development phase')
hline(sl, 1.35)

body(sl, [
    (0, 'Real ML Model Integration'),
    (1, 'Replace rule-based engine with TensorFlow/scikit-learn trained on MIMIC-III dataset'),
    (1, 'Model serving via TFServing sidecar container'),
    (0, 'FHIR R4 Interoperability'),
    (1, 'Expose FHIR-compliant REST endpoints for EHR system integration'),
    (1, 'Patient data export in HL7 FHIR Bundle format'),
    (0, 'Multi-Cluster Federation'),
    (1, 'Deploy across geo-distributed K3s clusters for HA and disaster recovery'),
    (0, 'Wearable Device Integration'),
    (1, 'Ingest Apple HealthKit / Google Fit data via MQTT + Kafka pipeline'),
    (0, 'Advanced Alerting'),
    (1, 'Grafana OnCall + PagerDuty integration for 24/7 clinical escalation'),
    (0, 'Cost Optimisation'),
    (1, 'Spot-instance node pools + Vertical Pod Autoscaler for right-sizing'),
], size=12)
pagenum(sl, 20)
print('Slide 20: Future Work')

# =============================================================================
# SLIDE 21 — THANK YOU / Q&A
# =============================================================================
sl = slide()
txt(sl, 'Thank You', 0, 1.8, 13.333, 1.1, size=50, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'Questions & Discussion', 0, 2.9, 13.333, 0.6, size=22, align=PP_ALIGN.CENTER)
hline(sl, 3.7)
txt(sl, 'Arti Kumari  |  2024MT03117  |  M.Tech Cloud Computing',
    0, 3.85, 13.333, 0.4, size=14, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'BITS Pilani — Work Integrated Learning Programmes  |  May 2026',
    0, 4.28, 13.333, 0.35, size=11, align=PP_ALIGN.CENTER)
hline(sl, 4.8)

qa = [
    ('Why K3s?',               '512MB RAM, full K8s API, 60MB binary, no vendor lock-in'),
    ('How does HPA work?',     'Metrics Server → 15s poll → ceil(actual/target) → 120s cooldown'),
    ('Why MySQL over NoSQL?',  'Structured health records, ACID compliance, familiar SQL ecosystem'),
    ('Why bcrypt factor 10?',  '~120ms compute per hash = brute-force economically infeasible'),
    ('How Prometheus scrapes?','Pull model: GET /metrics every 15s, K8s annotation discovery'),
    ('HIPAA coverage?',        'All 7 §164.312 technical safeguards addressed in architecture'),
]
for i, (q, a) in enumerate(qa):
    row, col = i // 2, i % 2
    x = 0.5 + col * 6.6
    y = 5.0 + row * 0.62
    txt(sl, f'Q: {q}', x, y, 6.2, 0.28, size=9.5, bold=True)
    txt(sl, f'A: {a}', x, y+0.3, 6.2, 0.28, size=9)

pagenum(sl, 21)
print('Slide 21: Q&A / Thank You')

# =============================================================================
# SAVE
# =============================================================================
prs.save(OUT)
print(f'\n✅  Saved: {OUT}')
print(f'   Slides: {len(prs.slides)}')
