"""Build complete 21-slide HealthAI presentation matching existing PPT theme."""
import sys, os, copy
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

TEMPLATE = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\2024MT03117- ABSTRACT.pptx'
OUT      = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\HealthAI_Presentation_v2.pptx'
SCRN     = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\screenshots'
IMGS     = r'd:\M.tech\SEM IV\Healthcare Dashboard Development\dissertation_imgs'

# ── theme colors (matching 2024MT03117- ABSTRACT.pptx) ───────────────────────
RED    = RGBColor(0xD7, 0x24, 0x2A)   # primary accent #D7242A
DARK   = RGBColor(0x1A, 0x1A, 0x1A)   # near-black text #1A1A1A
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF9, 0xF9, 0xF9)   # card background #F9F9F9
MGRAY  = RGBColor(0x55, 0x55, 0x55)   # secondary text #555555
DGRAY  = RGBColor(0x33, 0x33, 0x33)   # body text #333333
BLUE   = RGBColor(0x2E, 0x9B, 0xBE)   # blue accent #2E9BBE
GREEN  = RGBColor(0x2E, 0x7D, 0x32)   # green accent #2E7D32
ORANGE = RGBColor(0xF5, 0xA6, 0x23)   # orange accent #F5A623
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)   # purple accent #7B1FA2
LBLUE  = RGBColor(0xEB, 0xF4, 0xFA)   # light blue bg #EBF4FA

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation(TEMPLATE)
prs.slide_width  = W
prs.slide_height = H

# Use blank layout
blank_layout = prs.slide_layouts[6]   # Blank

# ── helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, l, t, w, h, fill_rgb=None, border_rgb=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if border_rgb and border_pt:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h, size=18, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def multiline_txbox(slide, lines, l, t, w, h, size=14, color=DGRAY, spacing=1.15):
    """lines = list of (text, bold, color_override_or_None)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (txt, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col if col else color
    return box

def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        print(f'  WARN: image not found: {path}')
        return None
    if h:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))

def red_badge(slide, num_text, l, t, size=0.55):
    r = rect(slide, l, t, size, size, RED)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num_text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE

def footer(slide, num, total=21):
    rect(slide, 0, 7.3, 13.333, 0.03, RGBColor(0xDD,0xDD,0xDD))
    txbox(slide, 'HealthAI  |  AI-Powered Health Symptom Tracker  |  M.Tech Cloud Computing  |  BITS Pilani WILP',
          0.3, 7.33, 11.5, 0.2, size=7.5, color=MGRAY)
    txbox(slide, f'{num} / {total}', 12.5, 7.33, 0.7, 0.2, size=8,
          color=MGRAY, align=PP_ALIGN.RIGHT)

def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.333, 0.07, RED)
    rect(slide, 0, 0.07, 0.07, 1.25, RED)
    txbox(slide, title, 0.25, 0.12, 12.8, 0.65, size=26, bold=True, color=DARK)
    if subtitle:
        txbox(slide, subtitle, 0.25, 0.72, 12.8, 0.32, size=11.5, color=MGRAY)
    rect(slide, 0.25, 1.1, 13.0, 0.03, RGBColor(0xE8,0xE8,0xE8))

def divider(slide, t):
    rect(slide, 0.35, t, 12.6, 0.035, RED)

def bullet_block(slide, items, l, t, w, h, title=None, title_color=RED,
                 item_size=13, title_size=15):
    """items = list of (indent, text) where indent 0=bullet, 1=sub"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for (indent, text) in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.level = indent
        p.space_after = Pt(1)
        if indent == 0:
            p.bullet = True
        run = p.add_run()
        run.text = text
        run.font.size = Pt(item_size - indent*1.5)
        run.font.color.rgb = DGRAY if indent == 0 else MGRAY
    return box

def code_block(slide, code_text, l, t, w, h, font_size=9):
    r = rect(slide, l, t, w, h, RGBColor(0xF0,0xF0,0xF0))
    box = slide.shapes.add_textbox(Inches(l+0.1), Inches(t+0.08), Inches(w-0.2), Inches(h-0.1))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in code_text.split('\n'):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = 'Courier New'
        run.font.color.rgb = DARK

def simple_table(slide, headers, rows, l, t, w, col_widths=None, header_bg=DARK,
                 row_height=0.32, font_size=11):
    """Draw a simple table as shapes."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w/n_cols] * n_cols
    # header row
    x = l
    rect(slide, l, t, w, row_height, RED)
    for ci, hdr in enumerate(headers):
        txbox(slide, hdr, x+0.05, t+0.05, col_widths[ci]-0.1, row_height-0.05,
              size=font_size, bold=True, color=WHITE)
        x += col_widths[ci]
    # data rows
    for ri, row in enumerate(rows):
        y = t + row_height*(ri+1)
        bg = LGRAY if ri%2==0 else WHITE
        rect(slide, l, y, w, row_height, bg)
        # vertical border
        rect(slide, l, y, w, row_height, border_rgb=RGBColor(0xCC,0xCC,0xCC), border_pt=0.5)
        x = l
        for ci, cell in enumerate(row):
            col = GREEN if cell in ('✅','✓','Yes') else (RED if cell in ('❌','No') else DGRAY)
            txbox(slide, cell, x+0.05, y+0.04, col_widths[ci]-0.1, row_height-0.05,
                  size=font_size-1, color=col)
            x += col_widths[ci]

# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
sl = add_slide()
# White background — no fill rect needed (default is white)
# Top red banner
rect(sl, 0, 0, 13.333, 1.6, RED)
txbox(sl, 'HealthAI', 0, 0.2, 13.333, 0.9, size=48, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, 'AI-Powered Health Symptom Tracker', 0, 1.05, 13.333, 0.45,
      size=16, color=RGBColor(0xFF,0xCC,0xCC), align=PP_ALIGN.CENTER)
# Center content on white background
txbox(sl, 'Deployed on Cloud-Native Infrastructure using Docker & Kubernetes (K3s)',
      1.5, 1.9, 10.3, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)
# Three pillars
pillars = [
    ('Patient Portal',    'Symptom logging, AI analysis,\nhealth history dashboard', BLUE),
    ('Admin Command',     'Real-time monitoring, alerts,\nsevere case management', RED),
    ('Cloud-Native Infra','Docker + K3s autoscaling,\nPrometheus + Grafana', GREEN),
]
for i, (title, body, col) in enumerate(pillars):
    x = 0.9 + i*4.15
    rect(sl, x, 2.55, 3.8, 2.4, LGRAY)
    rect(sl, x, 2.55, 3.8, 0.07, col)
    txbox(sl, title, x+0.15, 2.7, 3.5, 0.45, size=14, bold=True, color=col)
    txbox(sl, body, x+0.15, 3.2, 3.5, 0.8, size=11, color=DGRAY)
# Divider
rect(sl, 3.5, 5.2, 6.3, 0.04, RED)
# Author info
txbox(sl, 'Arti Kumari  |  2024MT03117  |  M.Tech Cloud Computing',
      0, 5.35, 13.333, 0.38, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
txbox(sl, 'BITS Pilani — Work Integrated Learning Programmes  |  May 2026',
      0, 5.72, 13.333, 0.35, size=11, color=MGRAY, align=PP_ALIGN.CENTER)
# Bottom red bar
rect(sl, 0, 7.3, 13.333, 0.2, RED)
txbox(sl, '1 / 21', 12.3, 7.3, 0.9, 0.18, size=9,
      color=WHITE, align=PP_ALIGN.RIGHT)
print('Slide 1: Title')

# =============================================================================
# SLIDE 2 — THE PROBLEM
# =============================================================================
sl = add_slide()
slide_header(sl, 'Why This Matters', 'The Problem Space in Digital Healthcare')
rect(sl, 0, 0, 0.06, 7.5, RED)

problems = [
    ('🔴 Reactive, not Proactive',
     'Patients log symptoms but get no AI-driven insights or early warnings',
     RED),
    ('🟡 Monolithic, not Scalable',
     'Single-server apps collapse under traffic spikes — no elasticity',
     RGBColor(0xE6,0x7E,0x00)),
    ('🟠 Observable on Paper, Blind in Practice',
     'No real-time dashboards — IT teams discover outages from angry users',
     RGBColor(0xCC,0x55,0x00)),
]
for i, (title, body, col) in enumerate(problems):
    x = 0.35 + i*4.3
    rect(sl, x, 1.45, 4.0, 2.8, LGRAY)
    rect(sl, x, 1.45, 4.0, 0.06, col)
    txbox(sl, title, x+0.15, 1.6, 3.7, 0.5, size=13, bold=True, color=col)
    txbox(sl, body, x+0.15, 2.1, 3.7, 1.0, size=11, color=DGRAY)

txbox(sl, 'The Stakes:', 0.35, 4.45, 3, 0.35, size=14, bold=True, color=DARK)
stakes = [
    '1 in 3 patients: "I didn\'t know when to see a doctor"',
    'Average healthcare app downtime: 14 hours / year',
    'HIPAA violations: avg fine $1.19M per incident (2023)',
]
for i, s in enumerate(stakes):
    txbox(sl, f'• {s}', 0.35, 4.85+i*0.38, 12.6, 0.35, size=12, color=DGRAY)
footer(sl, 2)
print('Slide 2: Problem')

# =============================================================================
# SLIDE 3 — SOLUTION
# =============================================================================
sl = add_slide()
slide_header(sl, 'Introducing HealthAI', 'One platform. Three capabilities.')
rect(sl, 0, 0, 0.06, 7.5, RED)

cols = [
    ('Patient Interface', ['Log symptoms in < 1 min', 'AI confidence scoring', 'Health trend charts', '6-category symptom grid'], '02_user_dashboard.png', BLUE),
    ('Admin Interface',   ['Monitor all patients', 'Severe case alerts', 'Population analytics', 'User growth charts'], '06_admin_dashboard.png', RED),
    ('Cloud Infrastructure', ['Auto-scale 2 → 10 pods', 'Prometheus + Grafana', 'HIPAA §164.312 compliant', '0.08% error @ 100 VUs'], 'fig4_1_docker.png', GREEN),
]
for i, (title, points, img, col) in enumerate(cols):
    x = 0.3 + i*4.35
    rect(sl, x, 1.45, 4.1, 5.5, LGRAY)
    rect(sl, x, 1.45, 4.1, 0.08, col)
    txbox(sl, title, x+0.1, 1.6, 3.9, 0.4, size=14, bold=True, color=col)
    for j, pt in enumerate(points):
        txbox(sl, f'✓  {pt}', x+0.15, 2.1+j*0.38, 3.8, 0.35, size=11, color=DGRAY)
    imgp = os.path.join(SCRN if not img.startswith('fig') else IMGS, img)
    if os.path.exists(imgp):
        add_image(sl, imgp, x+0.1, 3.75, 3.9, 2.9)

txbox(sl, '8 pages  •  15+ REST endpoints  •  6 custom Prometheus metrics  •  All HIPAA §164.312 safeguards covered',
      0.35, 7.0, 12.6, 0.3, size=10, color=MGRAY, align=PP_ALIGN.CENTER)
footer(sl, 3)
print('Slide 3: Solution')

# =============================================================================
# SLIDE 4 — TECHNOLOGY STACK
# =============================================================================
sl = add_slide()
slide_header(sl, 'Built With Production-Grade Tools', 'Full-stack technology selection')
rect(sl, 0, 0, 0.06, 7.5, RED)

stacks = [
    ('Frontend', [
        'React 18 + TypeScript — type-safe component UI',
        'Vite — sub-second builds, Hot Module Replacement',
        'Tailwind CSS v4 — utility-first responsive styling',
        'Recharts — LineChart, AreaChart, BarChart, RadarChart',
        'React Router v7 — nested layout routes',
        'shadcn/ui — accessible component library',
    ], BLUE),
    ('Backend', [
        'Node.js + Express.js — non-blocking REST API',
        'MySQL 8.x — relational health data storage',
        'bcryptjs (factor 10) — secure password hashing',
        'ioredis — Redis caching (30s TTL per user)',
        'express-rate-limit — brute-force protection',
        'prom-client — custom Prometheus metrics',
    ], RED),
    ('Infrastructure', [
        'Docker (multi-stage builds) — immutable containers',
        'Docker Compose — local orchestration',
        'K3s — lightweight Kubernetes (512MB RAM)',
        'Kubernetes HPA — CPU-based autoscaling (2–6 pods)',
        'Traefik — ingress controller + TLS termination',
        'Prometheus + Grafana — observability stack',
    ], GREEN),
]
for i, (title, items, col) in enumerate(stacks):
    x = 0.3 + i*4.35
    rect(sl, x, 1.45, 4.1, 5.5, LGRAY)
    rect(sl, x, 1.45, 0.06, 5.5, col)
    txbox(sl, title, x+0.2, 1.5, 3.8, 0.45, size=15, bold=True, color=col)
    divider(sl, 2.0) if i == 0 else None
    for j, item in enumerate(items):
        txbox(sl, f'• {item}', x+0.2, 2.05+j*0.58, 3.8, 0.52, size=10.5, color=DGRAY)
footer(sl, 4)
print('Slide 4: Tech Stack')

# =============================================================================
# SLIDE 5 — SYSTEM ARCHITECTURE
# =============================================================================
sl = add_slide()
slide_header(sl, 'Three-Tier Cloud-Native Architecture', 'System design overview')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(IMGS, 'fig3_1_arch.png'), 6.8, 1.4, 6.3, 5.5)

arch_text = [
    (0, 'Presentation Tier'),
    (1, 'React 18 + TypeScript SPA'),
    (1, 'Served by Nginx container'),
    (1, 'X-User-Id / X-User-Email headers'),
    (0, 'Application Tier'),
    (1, 'Node.js + Express.js REST API'),
    (1, 'Business logic + RBAC enforcement'),
    (1, 'AI analysis engine + prom-client'),
    (0, 'Data Tier'),
    (1, 'MySQL 8.x relational database'),
    (1, 'Redis cache (30s TTL)'),
    (0, 'Infrastructure Layer'),
    (1, 'Docker containers per service'),
    (1, 'K3s orchestration + Traefik ingress'),
    (1, 'Prometheus + Grafana monitoring'),
]
bullet_block(sl, arch_text, 0.35, 1.45, 6.2, 5.4, item_size=11.5)
footer(sl, 5)
print('Slide 5: Architecture')

# =============================================================================
# SLIDE 6 — FRONTEND ARCHITECTURE
# =============================================================================
sl = add_slide()
slide_header(sl, 'React 18 SPA — Role-Based Route Layout', 'Frontend structure and navigation')
rect(sl, 0, 0, 0.06, 7.5, RED)

code_block(sl, """/                    → RootLayout (public Navbar)
├── /                → HomePage
├── /login           → LoginPage
├── /signup          → SignupPage
└── /forgot-password → ForgotPasswordPage

/user                → UserLayout (sidebar, 5 nav items)
├── /user            → DashboardPage
├── /user/add-symptoms  → AddSymptomsPage
├── /user/history    → HistoryPage
├── /user/ai-suggestions → AISuggestionPage
└── /user/profile    → ProfilePage

/admin               → AdminLayout (admin sidebar)
├── /admin           → AdminDashboardPage
├── /admin/monitoring   → UserMonitoringPage
├── /admin/severe-cases → SevereCasesPage
├── /admin/alerts    → AlertsPage
└── /admin/reports   → ReportsPage""", 0.35, 1.45, 5.8, 5.5, font_size=9)

features = [
    ('Sidebar Features:', None),
    ('HealthAI logo + heart icon', None),
    ('User info card with live health score bar', None),
    ('Active symptom count badge on bell icon', None),
    ('HIPAA Compliant badge', None),
    ('Role-based navigation (patient vs admin)', None),
    ('', None),
    ('State Management:', None),
    ('localStorage for user session (id, email, role)', None),
    ('X-User-Id / X-User-Email headers on all API calls', None),
    ('No server-side sessions → enables HPA scaling', None),
]
y = 1.5
for (text, _) in features:
    bold = text.endswith(':')
    col = RED if bold else DGRAY
    txbox(sl, text, 6.4, y, 6.7, 0.35, size=11.5 if bold else 11,
          bold=bold, color=col)
    y += 0.37
footer(sl, 6)
print('Slide 6: Frontend')

# =============================================================================
# SLIDE 7 — PATIENT DASHBOARD
# =============================================================================
sl = add_slide()
slide_header(sl, 'Real-Time Patient Dashboard', 'What patients see on login')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(SCRN, '02_user_dashboard.png'), 6.8, 1.42, 6.3, 5.5)

features = [
    ('4 Live Metric Cards:', True, RED),
    ('Health Score — real-time with "Live" pulse', False, DGRAY),
    ('Active Symptoms count', False, DGRAY),
    ('AI Suggestions count', False, DGRAY),
    ('', False, DGRAY),
    ('Health Trend Chart (Recharts LineChart):', True, RED),
    ('7-day score trend with gradient fill', False, DGRAY),
    ('Animated tooltip on hover', False, DGRAY),
    ('', False, DGRAY),
    ('Microservices Status Panel:', True, RED),
    ('Real-time service health with response times', False, DGRAY),
    ('Animated green pulse per live service', False, DGRAY),
    ('', False, DGRAY),
    ('Recent Symptoms + AI Health Insights:', True, RED),
    ('Severity-colour-coded symptom list', False, DGRAY),
    ('High / Medium / Low priority AI badges', False, DGRAY),
]
multiline_txbox(sl, features, 0.35, 1.45, 6.2, 5.5, size=11.5)
footer(sl, 7)
print('Slide 7: Dashboard')

# =============================================================================
# SLIDE 8 — AI SYMPTOM ANALYSIS
# =============================================================================
sl = add_slide()
slide_header(sl, 'AI-Powered Symptom Analysis', 'How patients log and interpret symptoms')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(SCRN, '03_add_symptoms.png'), 0.3, 1.42, 5.8, 2.65)
add_image(sl, os.path.join(SCRN, '05_ai_suggestions.png'), 0.3, 4.15, 5.8, 2.65)

txbox(sl, 'Symptom Entry (AddSymptomsPage)', 6.35, 1.45, 6.7, 0.4, size=13, bold=True, color=RED)
steps = [
    '1. Select Category Tab — 6 types: Respiratory, Digestive,',
    '    Neurological, Musculoskeletal, Cardiovascular, General',
    '2. Toggle Symptoms — grid buttons per category',
    '3. Set Severity (1–10) — gradient slider + emoji scale',
    '4. Optional Temperature — fever alert if > 100.4°F',
    '5. AI Preview Card — preliminary analysis before submit',
    '6. Submit → POST /api/symptoms',
]
for i, s in enumerate(steps):
    txbox(sl, s, 6.35, 1.88+i*0.38, 6.7, 0.36, size=10.5, color=DGRAY)

txbox(sl, 'AI Insights (AISuggestionPage)', 6.35, 4.15, 6.7, 0.4, size=13, bold=True, color=RED)
ai_pts = [
    'Confidence Banner — animated % bar (green/yellow/red)',
    'Condition Accordion — expandable with causes & treatments',
    '"When to see a doctor" trigger warnings',
    'Priority-badged personalised recommendations',
    'Thumbs Up/Down feedback for model improvement',
    'Sparkline health trend at bottom',
]
for i, s in enumerate(ai_pts):
    txbox(sl, f'• {s}', 6.35, 4.6+i*0.4, 6.7, 0.38, size=10.5, color=DGRAY)
footer(sl, 8)
print('Slide 8: AI Analysis')

# =============================================================================
# SLIDE 9 — ADMIN INTERFACE
# =============================================================================
sl = add_slide()
slide_header(sl, 'Administrator Interface', 'Healthcare admin capabilities')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(SCRN, '06_admin_dashboard.png'), 0.3, 1.42, 6.2, 2.65)
add_image(sl, os.path.join(SCRN, '07_user_monitoring.png'), 0.3, 4.15, 6.2, 2.65)

panels = [
    ('Admin Dashboard', [
        'KPIs: Total Patients, Severe Cases, AI Queries, Uptime',
        '6-month User Growth AreaChart',
        'Hourly AI Query Volume BarChart',
        'Microservices health table: req/min, response time, uptime%',
        'Recent Activity Feed',
    ]),
    ('User Monitoring + Severe Cases', [
        'Searchable patient registry — filter by risk / status / name',
        'Server-side filtering: GET /api/admin/users?risk=high',
        'Row action: View Patient / Send Alert',
        'Severity ≥ 8 → auto-flagged as severe case',
        'One-click personalised alert sending',
    ]),
]
for i, (title, pts) in enumerate(panels):
    y = 1.45 + i*2.73
    txbox(sl, title, 6.75, y, 6.3, 0.4, size=13, bold=True, color=RED)
    for j, pt in enumerate(pts):
        txbox(sl, f'• {pt}', 6.75, y+0.45+j*0.42, 6.3, 0.4, size=10.5, color=DGRAY)
footer(sl, 9)
print('Slide 9: Admin')

# =============================================================================
# SLIDE 10 — DOCKER
# =============================================================================
sl = add_slide()
slide_header(sl, 'Multi-Stage Docker Builds', '97% frontend image size reduction')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(IMGS, 'fig4_1_docker.png'), 6.9, 1.42, 6.2, 3.5)

code_block(sl, """# Frontend — 850MB → 25MB (97% reduction)
FROM node:20-alpine AS builder
WORKDIR /app
COPY package*.json ./
RUN npm ci
COPY . .
RUN npm run build

FROM nginx:alpine AS production
COPY --from=builder /app/dist /usr/share/nginx/html
COPY nginx.conf /etc/nginx/conf.d/default.conf
EXPOSE 80""", 0.3, 1.45, 6.3, 2.85, font_size=9)

code_block(sl, """# Backend — 650MB → 180MB (72% reduction)
FROM node:20-alpine AS builder
WORKDIR /app
COPY package*.json ./
RUN npm ci --only=production

FROM node:20-alpine AS production
WORKDIR /app
COPY --from=builder /app/node_modules ./node_modules
COPY . .
EXPOSE 5000
CMD ["node", "server.js"]""", 0.3, 4.45, 6.3, 2.65, font_size=9)

simple_table(sl,
    ['Image', 'Before', 'After', 'Saving'],
    [['Frontend', '850 MB', '25 MB', '97%'],
     ['Backend',  '650 MB', '180 MB','72%']],
    l=6.9, t=5.1, w=6.2,
    col_widths=[2.1, 1.3, 1.3, 1.5], font_size=12)
txbox(sl, 'Benefits: Faster CI/CD  •  3.2s pod cold start (vs 18s+)  •  Smaller attack surface  •  Lower registry cost',
      0.3, 7.0, 13.0, 0.28, size=9.5, color=MGRAY, align=PP_ALIGN.CENTER)
footer(sl, 10)
print('Slide 10: Docker')

# =============================================================================
# SLIDE 11 — KUBERNETES (K3s)
# =============================================================================
sl = add_slide()
slide_header(sl, 'K3s — Kubernetes for the Real World', 'Lightweight production orchestration')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(IMGS, 'fig4_2_k3s.png'), 6.8, 1.42, 6.3, 3.8)

why = [
    ('Why K3s?', True, RED),
    ('Full Kubernetes API in a single ~60MB binary', False, DGRAY),
    ('Control plane: 512MB RAM (vs ~2GB for full K8s)', False, DGRAY),
    ('Built-in: Traefik ingress, CoreDNS, flannel CNI', False, DGRAY),
    ('Same YAML manifests → zero vendor lock-in', False, DGRAY),
    ('', False, DGRAY),
    ('HealthAI K8s Structure (namespace: healthcare)', True, RED),
    ('frontend/   Deployment (2 replicas) + Service', False, DGRAY),
    ('backend/    Deployment (2–6 pods, HPA) + Service', False, DGRAY),
    ('mysql/      Deployment + PVC + Service', False, DGRAY),
    ('redis/      Deployment + Service', False, DGRAY),
    ('monitoring/ Prometheus + Grafana + Loki', False, DGRAY),
    ('', False, DGRAY),
    ('Backend pod highlights:', True, RED),
    ('readinessProbe + livenessProbe on /api/health', False, DGRAY),
    ('Resource requests: CPU 100m–500m, RAM 128–512Mi', False, DGRAY),
    ('MySQL credentials from Kubernetes Secret', False, DGRAY),
    ('Redis host from ConfigMap (not hardcoded)', False, DGRAY),
]
multiline_txbox(sl, why, 0.35, 1.45, 6.2, 5.5, size=11)
footer(sl, 11)
print('Slide 11: K3s')

# =============================================================================
# SLIDE 12 — HPA
# =============================================================================
sl = add_slide()
slide_header(sl, 'Horizontal Pod Autoscaling', 'Elastic scaling under real load')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(IMGS, 'fig5_3_hpa.png'), 6.8, 1.42, 6.3, 3.5)

code_block(sl, """# Backend HPA — k8s/hpa.yaml
minReplicas: 2
maxReplicas: 6
metrics:
  - type: Resource
    resource:
      name: cpu
      target:
        type: Utilization
        averageUtilization: 60   # scale up if CPU > 60%
behavior:
  scaleUp:
    stabilizationWindowSeconds: 30   # fast scale-up
  scaleDown:
    stabilizationWindowSeconds: 120  # slow scale-down""", 0.3, 1.45, 6.3, 3.5, font_size=9)

txbox(sl, 'Observed Load Test Behaviour:', 0.3, 5.05, 6.3, 0.38, size=13, bold=True, color=RED)
simple_table(sl,
    ['Time', 'VUs', 'CPU%', 'Pods'],
    [['0:00', '0',   '12%', '2'],
     ['2:00', '100', '73%', '4  ↑'],
     ['3:30', '100', '68%', '6  ↑ (max)'],
     ['7:00', '0',   '8%',  '4  ↓'],
     ['15:00','0',   '5%',  '2  ↓']],
    l=0.3, t=5.48, w=6.3,
    col_widths=[1.3, 1.1, 1.3, 2.6], font_size=11, row_height=0.28)

txbox(sl, 'Result: 3× peak scaling  •  P95 < 200ms maintained  •  0.08% error rate at 100 VUs  ✅',
      6.8, 5.1, 6.3, 0.38, size=12, bold=True, color=GREEN)
add_image(sl, os.path.join(IMGS, 'fig5_1_response.png'), 6.8, 5.55, 6.3, 1.6)
footer(sl, 12)
print('Slide 12: HPA')

# =============================================================================
# SLIDE 13 — PROMETHEUS + GRAFANA
# =============================================================================
sl = add_slide()
slide_header(sl, 'Full-Stack Observability', 'Prometheus metrics + Grafana dashboards')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(SCRN, '11_grafana.png'), 6.8, 1.42, 6.3, 3.6)

simple_table(sl,
    ['Metric', 'Type', 'Shows'],
    [['http_requests_total',         'Counter',   'Request volume by route/status'],
     ['http_request_duration_seconds','Histogram', 'P50, P95, P99 latency'],
     ['symptoms_logged_total',       'Counter',   'Symptom submission rate'],
     ['ai_analysis_requests_total',  'Counter',   'AI query load'],
     ['ai_analysis_duration_seconds','Histogram', 'AI computation time'],
     ['auth_failures_total',         'Counter',   'Security monitoring']],
    l=0.3, t=1.45, w=6.3,
    col_widths=[2.5, 1.4, 2.4], font_size=9.5, row_height=0.3)

txbox(sl, 'Grafana Dashboard (4 rows):', 0.3, 4.4, 6.3, 0.38, size=13, bold=True, color=RED)
rows_info = [
    'Row 1: API Request Rate  •  Error Rate  •  P95 Latency  •  P99 Latency',
    'Row 2: Symptom Submission Rate  •  AI Query Rate  •  AI Duration P95',
    'Row 3: CPU per pod (stacked)  •  Memory per pod  •  HPA Replica Count',
    'Row 4: MySQL query duration  •  Redis hit/miss ratio',
]
for i, r in enumerate(rows_info):
    txbox(sl, f'• {r}', 0.3, 4.82+i*0.38, 6.3, 0.35, size=10, color=DGRAY)

txbox(sl, 'Alert Rules:', 6.8, 5.15, 6.3, 0.38, size=13, bold=True, color=RED)
alerts = [
    'HighErrorRate: >5% errors for 2 min → Warning',
    'HighLatency: P95 >500ms for 5 min → Warning',
    'BackendDown: pod unreachable for 1 min → Critical',
]
for i, a in enumerate(alerts):
    txbox(sl, f'• {a}', 6.8, 5.58+i*0.38, 6.3, 0.35, size=11, color=DGRAY)
footer(sl, 13)
print('Slide 13: Monitoring')

# =============================================================================
# SLIDE 14 — SECURITY
# =============================================================================
sl = add_slide()
slide_header(sl, 'Security by Design — Defence in Depth', '7-layer security architecture')
rect(sl, 0, 0, 0.06, 7.5, RED)

layers = [
    ('Layer 1: Network',      'TLS 1.3 — all external traffic encrypted in transit'),
    ('Layer 2: Gateway',      'Traefik Ingress — HTTPS enforcement, HTTP redirect'),
    ('Layer 3: Application',  'Helmet.js — CSP, X-Frame-Options, X-Content-Type-Options'),
    ('Layer 4: Auth',         'bcryptjs (factor 10) — 120ms hash, brute-force resistant'),
    ('Layer 5: Rate Limiting','10 auth requests / 15 min per IP (express-rate-limit)'),
    ('Layer 6: Access Control','RBAC — patients: /user only | admins: /admin + /user'),
    ('Layer 7: Data',         'AES-256 encryption for PHI + MySQL parameterised queries'),
]
for i, (layer, desc) in enumerate(layers):
    y = 1.5 + i*0.63
    rect(sl, 0.3, y, 12.7, 0.58, LGRAY if i%2==0 else WHITE)
    rect(sl, 0.3, y, 0.05, 0.58, RED)
    txbox(sl, layer, 0.5, y+0.05, 3.2, 0.48, size=11.5, bold=True, color=RED)
    txbox(sl, desc,  3.8, y+0.05, 9.0, 0.48, size=11.5, color=DGRAY)

txbox(sl, 'Auth Flow: POST /api/auth/login → { user: {id, name, email, role} } → All API calls send X-User-Id header → Backend RBAC check per route',
      0.3, 6.98, 12.7, 0.25, size=9, color=MGRAY, align=PP_ALIGN.CENTER)
footer(sl, 14)
print('Slide 14: Security')

# =============================================================================
# SLIDE 15 — HIPAA
# =============================================================================
sl = add_slide()
slide_header(sl, 'HIPAA §164.312 — Full Technical Safeguard Coverage', 'Compliance by architecture')
rect(sl, 0, 0, 0.06, 7.5, RED)

simple_table(sl,
    ['HIPAA Requirement', 'HealthAI Implementation', 'Status'],
    [['Access Control',             'RBAC: patient/admin roles enforced on every API route', '✅'],
     ['Audit Controls',             'Auth events + admin actions logged with timestamp + IP', '✅'],
     ['Integrity Controls',         'AES-256-GCM authentication tag on PHI at rest', '✅'],
     ['Transmission Security',      'TLS 1.3 mandatory — HTTP automatically redirected', '✅'],
     ['Person/Entity Authentication','bcryptjs + rate-limited login (10 req/15 min)', '✅'],
     ['Automatic Logoff',           'Session timeout on inactivity (frontend enforced)', '✅'],
     ['Encryption/Decryption',      'AES-256 at rest  +  TLS 1.3 in transit', '✅']],
    l=0.3, t=1.45, w=12.7,
    col_widths=[3.4, 6.8, 2.5], font_size=11.5, row_height=0.45)

txbox(sl, 'OWASP Top 10 Mitigations Applied:', 0.3, 5.2, 7, 0.38, size=13, bold=True, color=RED)
owasp = 'SQL Injection ✅  •  Broken Auth ✅  •  XSS ✅  •  IDOR ✅  •  Security Misconfiguration ✅  •  Sensitive Data Exposure ✅'
txbox(sl, owasp, 0.3, 5.6, 12.7, 0.4, size=11.5, color=DGRAY)

txbox(sl, 'Patient Profile shows: "AES-256 Encryption Active" + "TLS 1.3 Secured" + "HIPAA Compliant Storage"',
      0.3, 6.1, 12.7, 0.35, size=11, color=BLUE)
add_image(sl, os.path.join(SCRN, '14_user_profile.png'), 7.8, 5.18, 5.2, 1.9)
footer(sl, 15)
print('Slide 15: HIPAA')

# =============================================================================
# SLIDE 16 — LOAD TESTING
# =============================================================================
sl = add_slide()
slide_header(sl, 'System Validated Under Real Load', 'k6 load test results')
rect(sl, 0, 0, 0.06, 7.5, RED)

add_image(sl, os.path.join(IMGS, 'fig5_2_cpu_mem.png'), 0.3, 1.45, 6.0, 3.2)

txbox(sl, 'Test Setup (k6):', 6.5, 1.45, 6.5, 0.38, size=13, bold=True, color=RED)
setup = [
    '0 → 100 virtual users over 2 minutes',
    'Hold at 100 VUs for 3 minutes, then ramp down',
    'Request mix: 40% dashboard, 30% symptoms,',
    '20% AI analysis, 10% authentication',
]
for i, s in enumerate(setup):
    txbox(sl, f'• {s}', 6.5, 1.88+i*0.38, 6.5, 0.35, size=11, color=DGRAY)

simple_table(sl,
    ['Metric', 'Result', 'Target', '✓'],
    [['Total Requests',    '42,350',   '—',        '—'],
     ['Peak Rate',         '141 req/s','—',        '—'],
     ['P50 Response Time', '48 ms',    '—',        '✅'],
     ['P95 Response Time', '187 ms',   '< 200ms',  '✅'],
     ['P99 Response Time', '312 ms',   '—',        '✅'],
     ['Error Rate',        '0.08%',    '< 0.1%',   '✅'],
     ['Peak Pods',         '6 (from 2)','2–6',     '✅']],
    l=6.5, t=3.5, w=6.5,
    col_widths=[2.8, 1.4, 1.4, 0.9], font_size=11, row_height=0.3)

txbox(sl, 'HPA validated: 2 pods → 6 pods in response to 73% CPU — zero manual intervention',
      0.3, 4.78, 6.0, 0.45, size=11.5, bold=True, color=GREEN)
add_image(sl, os.path.join(IMGS, 'fig5_1_response.png'), 0.3, 5.3, 6.0, 1.85)
footer(sl, 16)
print('Slide 16: Load Testing')

# =============================================================================
# SLIDE 17 — PERFORMANCE
# =============================================================================
sl = add_slide()
slide_header(sl, 'Frontend Performance — Lighthouse Scores', 'Build optimisation results')
rect(sl, 0, 0, 0.06, 7.5, RED)

# Score circle (simulate)
rect(sl, 0.5, 1.6, 2.8, 2.8, LGRAY)
txbox(sl, '94', 0.5, 2.0, 2.8, 1.6, size=64, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(sl, '/100', 0.5, 3.1, 2.8, 0.5, size=16, color=MGRAY, align=PP_ALIGN.CENTER)
txbox(sl, 'Lighthouse\nPerformance', 0.5, 3.65, 2.8, 0.7, size=12, color=DARK,
      align=PP_ALIGN.CENTER)

simple_table(sl,
    ['Metric', 'Score', 'Grade'],
    [['Performance Score',      '94/100', 'A'],
     ['First Contentful Paint', '0.8 s',  'Excellent'],
     ['Largest Contentful Paint','1.2 s',  'Good'],
     ['Time to Interactive',    '1.4 s',  'Good'],
     ['Total Blocking Time',    '12 ms',  'Excellent'],
     ['Cumulative Layout Shift','0.02',   'Excellent']],
    l=3.5, t=1.45, w=6.0,
    col_widths=[3.4, 1.3, 1.3], font_size=12, row_height=0.37)

txbox(sl, 'Why scores are high:', 9.7, 1.45, 3.5, 0.38, size=13, bold=True, color=RED)
reasons = [
    'Nginx gzip compression (assets > 1KB)',
    'Cache-Control: immutable on hashed JS/CSS',
    'Vite code-splitting — loads only current page',
    '25MB Docker image → 3.2s pod cold start',
    'Cache-Control: no-cache on index.html',
]
for i, r in enumerate(reasons):
    txbox(sl, f'• {r}', 9.7, 1.9+i*0.45, 3.5, 0.42, size=10.5, color=DGRAY)

txbox(sl, 'Backend Performance (normal load, 0–20 VUs):', 0.3, 5.65, 12.7, 0.38,
      size=13, bold=True, color=RED)
simple_table(sl,
    ['Percentile', 'Response Time', 'SLO Status'],
    [['P50 (median)', '18 ms', 'Excellent — well within SLO'],
     ['P95',          '62 ms', 'Excellent'],
     ['P99',          '98 ms', 'Within 200ms SLO ✅']],
    l=0.3, t=6.1, w=9.0,
    col_widths=[2.4, 2.4, 4.2], font_size=12, row_height=0.3)
footer(sl, 17)
print('Slide 17: Performance')

# =============================================================================
# SLIDE 18 — SECURITY SCAN
# =============================================================================
sl = add_slide()
slide_header(sl, 'OWASP ZAP Scan — Zero High/Critical Findings', 'Security validation results')
rect(sl, 0, 0, 0.06, 7.5, RED)

simple_table(sl,
    ['Alert', 'Risk', 'Resolution'],
    [['Missing Anti-CSRF Token',        'Medium', 'Accepted — REST API uses CORS, not form tokens'],
     ['X-Content-Type-Options Missing', 'Low',    'Fixed — Helmet.js configured'],
     ['Content-Security-Policy Missing','Medium', 'Fixed — Helmet CSP configured'],
     ['Server Version Leakage',         'Low',    'Fixed — Express version header removed'],
     ['Insecure Form Action',           'Info',   'False Positive — React controlled inputs'],
     ['SQL Injection',                  'Info',   'N/A — MySQL parameterised queries used']],
    l=0.3, t=1.45, w=12.7,
    col_widths=[3.6, 1.5, 7.6], font_size=11, row_height=0.38)

txbox(sl, 'Manual Security Tests — All Passed:', 0.3, 4.25, 12.7, 0.38, size=13, bold=True, color=RED)
tests = [
    '✅  IDOR: User A cannot read User B\'s symptoms (confirmed via direct API test)',
    '✅  RBAC: Patient token → 403 Forbidden on all /api/admin/* routes (confirmed)',
    '✅  Rate Limiting: 11th auth request → 429 Too Many Requests (confirmed)',
    '✅  Password Policy: Minimum 8 characters enforced at signup (confirmed)',
    '✅  XSS: React JSX escaping prevents script injection (confirmed)',
]
for i, t in enumerate(tests):
    txbox(sl, t, 0.3, 4.7+i*0.42, 12.7, 0.4, size=11.5, color=DGRAY if i>0 else GREEN)

txbox(sl, 'Conclusion: No High or Critical vulnerabilities. All Medium findings remediated or accepted with documented rationale.',
      0.3, 6.88, 12.7, 0.28, size=10.5, bold=True, color=GREEN)
footer(sl, 18)
print('Slide 18: Security Scan')

# =============================================================================
# SLIDE 19 — CONCLUSIONS
# =============================================================================
sl = add_slide()
slide_header(sl, 'What We Proved', '6 key findings from this dissertation')
rect(sl, 0, 0, 0.06, 7.5, RED)

findings = [
    ('1', 'Full-Stack Healthcare SaaS is achievable',
     'React 18 + Vite + Tailwind CSS v4 → 8 production-quality pages in one cohesive system'),
    ('2', 'K3s delivers real Kubernetes at 512MB RAM',
     'Full API surface, identical YAML manifests — no infrastructure-size compromise'),
    ('3', 'HPA works — the numbers prove it',
     '2 → 6 pods automatically, P95 < 200ms maintained, 0.08% error rate at 100 VUs'),
    ('4', 'Security-by-design eliminates retrofitting',
     'All 7 HIPAA §164.312 safeguards addressed from architecture phase, not added later'),
    ('5', 'Observability enables confidence, not guesswork',
     'Prometheus + Grafana gave quantitative evidence of HPA behaviour and latency profile'),
    ('6', 'Multi-stage builds matter in practice',
     '97% image size reduction → faster CI/CD, faster K8s cold start, smaller attack surface'),
]
for i, (num, title, body) in enumerate(findings):
    row = i // 2
    col = i % 2
    x = 0.3 + col * 6.55
    y = 1.45 + row * 1.82
    rect(sl, x, y, 6.3, 1.72, LGRAY)
    red_badge(sl, num, x+0.1, y+0.08, size=0.55)
    txbox(sl, title, x+0.75, y+0.1, 5.4, 0.42, size=12.5, bold=True, color=RED)
    txbox(sl, body,  x+0.15, y+0.6,  6.0, 0.95, size=10.5, color=DGRAY)
footer(sl, 19)
print('Slide 19: Conclusions')

# =============================================================================
# SLIDE 20 — FUTURE WORK
# =============================================================================
sl = add_slide()
slide_header(sl, 'What Comes Next', 'Roadmap for future development')
rect(sl, 0, 0, 0.06, 7.5, RED)

simple_table(sl,
    ['Priority', 'Feature', 'Why'],
    [['High',   'HL7 FHIR API Integration',        'Connect patient data to clinical EHRs (Epic/Cerner)'],
     ['High',   'ML Symptom Analysis (BioMedBERT)', 'Replace rule-based engine with trained model'],
     ['Medium', 'Multi-Region K3s Deployment',      'Geographic redundancy + GDPR/PDPA compliance'],
     ['Medium', 'WebSocket Real-Time Notifications','Push alerts without page refresh'],
     ['Medium', 'React Native Mobile Apps',         'Wearable device integration + Bluetooth sensors'],
     ['Low',    'Service Mesh (Istio/Linkerd)',      'mTLS for pod-to-pod zero-trust communication'],
     ['Low',    'Chaos Engineering (Litmus)',        'Empirically validate resilience under failure']],
    l=0.3, t=1.45, w=12.7,
    col_widths=[1.6, 3.8, 7.3], font_size=11.5, row_height=0.42)

txbox(sl, 'Immediate Next Steps (if production):', 0.3, 5.55, 12.7, 0.38, size=13, bold=True, color=RED)
next_steps = [
    '1.  Optimise AI analysis endpoint (P95 was 342ms — add Redis caching with 60s TTL)',
    '2.  Add etcd cluster to K3s for control plane high availability',
    '3.  Complete GitHub Actions CI/CD pipeline — automated image build + kubectl rollout',
]
for i, s in enumerate(next_steps):
    txbox(sl, s, 0.3, 6.0+i*0.38, 12.7, 0.35, size=11.5, color=DGRAY)
footer(sl, 20)
print('Slide 20: Future Work')

# =============================================================================
# SLIDE 21 — Q&A / THANK YOU
# =============================================================================
sl = add_slide()
# Top red banner
rect(sl, 0, 0, 13.333, 2.2, RED)
txbox(sl, '♥', 0, 0.2, 13.333, 0.9, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, 'Thank You', 0, 1.05, 13.333, 0.85, size=38, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
# Author on white bg
txbox(sl, 'Arti Kumari  |  2024MT03117  |  M.Tech Cloud Computing',
      0, 2.35, 13.333, 0.45, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
txbox(sl, 'BITS Pilani — Work Integrated Learning Programmes  |  May 2026',
      0, 2.8, 13.333, 0.35, size=11, color=MGRAY, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 3.25, 6.3, 0.04, RED)

# Quick reference Q&A
qa = [
    ('Why K3s?',              '512MB RAM, full K8s API, single 60MB binary, no vendor lock-in'),
    ('How does HPA work?',    'Metrics Server → 15s poll → ceil(actual/target) → 120s cooldown'),
    ('Why MySQL?',            'Structured health records, ACID compliance, familiar SQL queries'),
    ('Why bcrypt factor 10?', '~120ms compute per hash = brute-force economically infeasible'),
    ('How Prometheus scrapes?','Pull model: GET /metrics every 15s, K8s annotation discovery'),
    ('HIPAA coverage?',       'All 7 §164.312 technical safeguards addressed in architecture'),
]
for i, (q, a) in enumerate(qa):
    row, col = i//2, i%2
    x = 0.5 + col*6.6
    y = 3.45 + row*0.62
    rect(sl, x, y, 6.2, 0.55, LGRAY)
    txbox(sl, f'Q: {q}', x+0.1, y+0.04, 6.0, 0.24, size=9.5, bold=True, color=RED)
    txbox(sl, f'A: {a}', x+0.1, y+0.28, 6.0, 0.24, size=9, color=DGRAY)

rect(sl, 0, 7.3, 13.333, 0.2, RED)
txbox(sl, 'HealthAI — Cloud-Native  |  Secure  |  Observable  |  HIPAA Compliant',
      0, 7.3, 13.0, 0.18, size=8, color=WHITE, align=PP_ALIGN.CENTER)
print('Slide 21: Q&A / Thank You')

# =============================================================================
# REMOVE ORIGINAL SLIDES (keep only new ones)
# =============================================================================
xml_slides = prs.slides._sldIdLst
# Remove the original 16 slides (they are at indices 0–15 in xml)
original_count = 16
for _ in range(original_count):
    rId = prs.slides._sldIdLst[0].get('r:id') if prs.slides._sldIdLst[0].get('r:id') else prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sl_elem = prs.slides._sldIdLst[0]
    prs.slides._sldIdLst.remove(sl_elem)

prs.save(OUT)
print(f'\n✅  Saved: {OUT}')
print(f'   Slides: {len(prs.slides)}')
