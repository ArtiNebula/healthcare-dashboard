import sys, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation(r'd:\M.tech\SEM IV\Healthcare Dashboard Development\2024MT03117- ABSTRACT.pptx')
master = prs.slide_master

theme_xml = master.element.xml
pat = r'val="([0-9A-Fa-f]{6})"'
colors = re.findall(pat, theme_xml)
print('Colors in master:', list(dict.fromkeys(colors))[:20])

# Layout names
for i, layout in enumerate(prs.slide_layouts):
    print(f'Layout {i}: {layout.name}')

# Check slides content
for idx, slide in enumerate(prs.slides):
    print(f'\n--- Slide {idx+1} ---')
    for shape in slide.shapes:
        try:
            if shape.text.strip():
                print(f'  text: {repr(shape.text[:80])}')
        except: pass
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    rgb = fill.fore_color.rgb
                    print(f'  fill: #{rgb}')
                except: pass
        except: pass
